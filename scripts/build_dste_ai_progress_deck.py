#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成《DSTE + AI 建设进展汇报》PPT。

基于帆软风格设计规范，从空白模板构建 17 页内部进展汇报 PPT。
输出: /Users/jasonjing/DSTE+AI内部进展汇报-v2.pptx

用法:
    .pptx-venv/bin/python scripts/build_dste_ai_progress_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path("/Users/jasonjing/DSTE+AI内部进展汇报-v2.pptx")
FONT = "PingFang SC"

# 帆软风格配色
BLUE = RGBColor(0x03, 0x5D, 0xCF)        # 主色
BLUE_MID = RGBColor(0x5B, 0x8C, 0xC8)    # 中蓝
BLUE_LIGHT = RGBColor(0xE8, 0xF4, 0xFD)  # 浅蓝
BLUE_PALE = RGBColor(0xF0, 0xF7, 0xFF)   # 极浅蓝
INK = RGBColor(0x17, 0x34, 0x6E)         # 深蓝正文
TEXT = RGBColor(0x33, 0x33, 0x33)        # 正文
MUTED = RGBColor(0x90, 0x98, 0xAA)       # 辅助灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0xB5, 0x78)       # 进展/达成
ORANGE = RGBColor(0xFF, 0x6A, 0x00)      # 预警/强调
DARK_BG = RGBColor(0x16, 0x1B, 0x26)     # 深色底座


def set_text(tf, runs_spec, align=PP_ALIGN.LEFT, space_after=Pt(4), line_spacing=1.0):
    """runs_spec: list of paragraphs; each paragraph is list of (text, size, bold, color)."""
    tf.word_wrap = True
    for i, para in enumerate(runs_spec):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.line_spacing = line_spacing
        for text, size, bold, color in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def add_box(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=False):
    box = slide.shapes.add_shape(shape, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if not line:
        box.line.fill.background()
    box.shadow.inherit = False
    return box


def add_textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.vertical_anchor = valign
    set_text(tf, runs, align=align, line_spacing=line_spacing)
    return tb


def add_section_title(slide, title, subtitle=None):
    """页面统一标题区。"""
    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.6),
                [[(title, 24, True, BLUE)]])
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.35),
                    [[(subtitle, 12, False, MUTED)]])


def add_footer(slide, text="DSTE + AI 建设进展汇报 · v0.7.14 · 2026-07-29"):
    add_textbox(slide, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.25),
                [[(text, 9, False, MUTED)]])


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景装饰
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_box(slide, Inches(8.5), Inches(0), Inches(4.83), prs.slide_height, BLUE_PALE, MSO_SHAPE.RECTANGLE)
    # 大标题
    add_textbox(slide, Inches(0.9), Inches(2.2), Inches(10.5), Inches(1.2),
                [[("DSTE + AI", 44, True, BLUE), (" 建设进展汇报", 44, True, INK)]])
    # 副标题
    add_textbox(slide, Inches(0.9), Inches(3.5), Inches(10), Inches(0.7),
                [[("以 AI 重构战略执行全链路，从「记录工具」迈向「智能伙伴」", 16, False, MUTED)]])
    # 关键信息卡片
    add_box(slide, Inches(0.9), Inches(4.6), Inches(3.2), Inches(1.1), BLUE_LIGHT)
    add_textbox(slide, Inches(1.05), Inches(4.75), Inches(2.9), Inches(0.8),
                [[("当前版本", 11, False, MUTED)],
                 [("v0.7.14", 22, True, BLUE)]], valign=MSO_ANCHOR.MIDDLE)
    add_box(slide, Inches(4.35), Inches(4.6), Inches(3.2), Inches(1.1), BLUE_LIGHT)
    add_textbox(slide, Inches(4.5), Inches(4.75), Inches(2.9), Inches(0.8),
                [[("汇报日期", 11, False, MUTED)],
                 [("2026-07-29", 22, True, BLUE)]], valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    return slide


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, "今天汇报什么")
    items = [
        "背景与目标",
        "建设进展",
        "产品架构",
        "AI 能力地图",
        "前后对比",
        "核心场景",
        "最新亮点",
        "成果数据",
        "下一步计划",
    ]
    # 垂直步骤条
    x, y = 1.5, 1.7
    for i, item in enumerate(items, 1):
        # 圆圈序号
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y + (i - 1) * 0.58), Inches(0.38), Inches(0.38))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE if i == 1 else BLUE_LIGHT
        circle.line.fill.background()
        add_textbox(slide, Inches(x), Inches(y + (i - 1) * 0.58), Inches(0.38), Inches(0.38),
                    [[(str(i), 12, True, WHITE if i == 1 else BLUE)]], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # 文字
        add_textbox(slide, Inches(x + 0.55), Inches(y + (i - 1) * 0.58 + 0.05), Inches(6), Inches(0.3),
                    [[(item, 14, i == 1, INK if i == 1 else MUTED)]])
        if i < len(items):
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.17), Inches(y + (i - 1) * 0.58 + 0.4), Inches(0.04), Inches(0.18))
            line.fill.solid()
            line.fill.fore_color.rgb = BLUE_LIGHT
            line.line.fill.background()
    add_footer(slide)


def slide_pain_points(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, "战略执行数字化面临的现实挑战")
    pains = [
        ("战略解码难落地", "年度目标、战略地图拆解后，与执行动作关联弱"),
        ("目标与执行脱节", "RoadMap、行动计划、经分会决议分散在不同模块"),
        ("数据分散不统一", "经营数据、填报数据、外部系统数据未形成统一视图"),
        ("协同成本高", "跨部门目标依赖、资源协调、会议组织大量依赖人工"),
        ("经验难以复用", "历史决议、最佳实践、风险案例未沉淀"),
        ("风险发现滞后", "指标偏离、行动项逾期多在复盘时才暴露"),
    ]
    # 2x3 卡片网格
    cols, rows = 3, 2
    margin_x, margin_y = 0.7, 1.7
    gap_x, gap_y = 0.3, 0.3
    card_w = (13.333 - 2 * margin_x - (cols - 1) * gap_x) / cols
    card_h = (7.5 - margin_y - 1.0 - (rows - 1) * gap_y) / rows
    for idx, (title, desc) in enumerate(pains):
        c, r = idx % cols, idx // cols
        x = margin_x + c * (card_w + gap_x)
        y = margin_y + r * (card_h + gap_y)
        add_box(slide, Inches(x), Inches(y), Inches(card_w), Inches(card_h), BLUE_LIGHT)
        # 左侧色条
        add_box(slide, Inches(x), Inches(y), Inches(0.06), Inches(card_h), BLUE, MSO_SHAPE.RECTANGLE)
        add_textbox(slide, Inches(x + 0.18), Inches(y + 0.15), Inches(card_w - 0.3), Inches(0.4),
                    [[(title, 13, True, INK)]])
        add_textbox(slide, Inches(x + 0.18), Inches(y + 0.6), Inches(card_w - 0.3), Inches(card_h - 0.75),
                    [[(desc, 10.5, False, MUTED)]], line_spacing=1.15)
    add_footer(slide)


def slide_value_positioning(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, "DSTE + AI 要达成什么")
    # 金句
    add_box(slide, Inches(0.7), Inches(1.35), Inches(11.9), Inches(0.8), BLUE_LIGHT)
    add_textbox(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
                [[("让 DSTE 从「战略执行记录工具」升级为「战略执行智能伙伴」", 16, True, INK)]],
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    # 四象限
    values = [
        ("战略解码智能化", "从人工拆解到 AI 辅助生成", BLUE_LIGHT),
        ("执行跟踪实时化", "从定期汇报到持续感知", BLUE_PALE),
        ("经营分析高效化", "从流程在线到决策智能", BLUE_PALE),
        ("组织知识资产化", "从信息孤岛到语义复用", BLUE_LIGHT),
    ]
    center_x, center_y = 6.666, 4.35
    quad_w, quad_h = 5.5, 2.0
    positions = [(-1, -1), (1, -1), (-1, 1), (1, 1)]
    for (title, desc, fill), (dx, dy) in zip(values, positions):
        x = center_x + dx * (quad_w / 2 + 0.1)
        y = center_y + dy * (quad_h / 2 + 0.1)
        add_box(slide, Inches(x - quad_w / 2), Inches(y - quad_h / 2), Inches(quad_w), Inches(quad_h), fill)
        add_textbox(slide, Inches(x - quad_w / 2 + 0.15), Inches(y - quad_h / 2 + 0.15), Inches(quad_w - 0.3), Inches(0.45),
                    [[(title, 13, True, BLUE)]], align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x - quad_w / 2 + 0.15), Inches(y - quad_h / 2 + 0.6), Inches(quad_w - 0.3), Inches(0.8),
                    [[(desc, 10.5, False, MUTED)]], align=PP_ALIGN.CENTER, line_spacing=1.1)
    # 中心圆
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - 0.55), Inches(center_y - 0.55), Inches(1.1), Inches(1.1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    add_textbox(slide, Inches(center_x - 0.55), Inches(center_y - 0.55), Inches(1.1), Inches(1.1),
                [[("AI", 16, True, WHITE)]], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)


def slide_milestones(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, "从 v0.6 到 v0.7，能力快速收敛")
    milestones = [
        ("v0.6.0–v0.6.5", "人员/组织目录、战略洞察与专题拆分、OMP 扩展、AI 助手迭代"),
        ("v0.6.6–v0.6.11", "设计系统 Phase 1、OMP 云端同步、全模块 401 收尾、同步健壮性"),
        ("v0.6.12–v0.6.17", "年度计划隔离、AI 抽屉右侧化、KMS AI 问答、行动项 progressLogs"),
        ("v0.7.0–v0.7.14", "规则引擎中心、驾驶舱多标签工作区、全局待办面板、目录管理"),
    ]
    # 横向时间轴
    y = 2.4
    total_w = 11.9
    n = len(milestones)
    step = total_w / (n - 1)
    # 主线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y + 0.6), Inches(total_w), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE_LIGHT
    line.line.fill.background()
    for i, (ver, desc) in enumerate(milestones):
        x = 0.7 + i * step
        # 节点
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.12), Inches(y + 0.6 - 0.12), Inches(0.24), Inches(0.24))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE if i == n - 1 else BLUE_MID
        circle.line.fill.background()
        # 版本号
        add_textbox(slide, Inches(x - 0.9), Inches(y - 0.1), Inches(1.8), Inches(0.35),
                    [[(ver, 12, True, BLUE if i == n - 1 else INK)]], align=PP_ALIGN.CENTER)
        # 描述卡片
        card_y = y + 1.1
        add_box(slide, Inches(x - 1.05), Inches(card_y), Inches(2.1), Inches(1.8), BLUE_LIGHT if i < n - 1 else BLUE)
        add_textbox(slide, Inches(x - 0.95), Inches(card_y + 0.12), Inches(1.9), Inches(1.6),
                    [[(desc, 9.5, False, TEXT if i < n - 1 else WHITE)]], line_spacing=1.1)
    # 当前版本标签
    add_box(slide, Inches(10.8), Inches(6.3), Inches(1.8), Inches(0.5), ORANGE)
    add_textbox(slide, Inches(10.8), Inches(6.3), Inches(1.8), Inches(0.5),
                [[("v0.7.14", 12, True, WHITE)]], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, "DSTE + AI 技术架构：四层一体")
    layers = [
        ("交互层", "Web 端 / 移动端 / 会议大屏 / 企业微信", BLUE_PALE),
        ("AI 能力层", "NLP · 推荐引擎 · 预测预警 · Agent 编排", BLUE_LIGHT),
        ("业务中台", "会议中心 / 决议中心 / RoadMap / 战略地图 / 人员组织 / 规则引擎", BLUE_PALE),
        ("数据底座", "填报数据 + 业务系统数据 + KMS 知识库 / 统一持久化 / 向量检索 / 版本审计", BLUE_LIGHT),
    ]
    x, y, w, h = 0.9, 1.7, 11.5, 4.8
    layer_h = h / len(layers)
    for i, (title, desc, fill) in enumerate(layers):
        ly = y + i * layer_h
        add_box(slide, Inches(x), Inches(ly), Inches(w), Inches(layer_h - 0.08), fill)
        add_textbox(slide, Inches(x + 0.2), Inches(ly + 0.12), Inches(2.2), Inches(0.4),
                    [[(title, 14, True, BLUE)]])
        add_textbox(slide, Inches(x + 2.5), Inches(ly + 0.18), Inches(w - 2.7), Inches(0.35),
                    [[(desc, 11, False, INK)]])
    # 右侧技术标签
    add_box(slide, Inches(10.5), Inches(1.85), Inches(2.2), Inches(0.45), BLUE)
    add_textbox(slide, Inches(10.5), Inches(1.85), Inches(2.2), Inches(0.45),
                [[("Kimi 大模型", 10, True, WHITE)]], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_box(slide, Inches(10.5), Inches(2.45), Inches(2.2), Inches(0.45), BLUE)
    add_textbox(slide, Inches(10.5), Inches(2.45), Inches(2.2), Inches(0.45),
                [[("SSE 流式", 10, True, WHITE)]], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_box(slide, Inches(10.5), Inches(3.05), Inches(2.2), Inches(0.45), BLUE)
    add_textbox(slide, Inches(10.5), Inches(3.05), Inches(2.2), Inches(0.45),
                [[("RAG + 向量库", 9, True, WHITE)]], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)


def slide_data_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, "数据架构：以服务器为唯一真相源")
    layers = [
        ("接入层", "填报数据、RoadMap、人员组织、外部业务系统、KMS", "统一接入"),
        ("治理层", "数据清洗、实体关联、权限脱敏、版本管理", "保证质量"),
        ("知识层", "会议知识库、决议库、指标库、最佳实践库", "沉淀复用"),
        ("服务层", "标准 API / 向量检索 / 实时数据推送", "支撑 AI 与业务"),
    ]
    x, y = 0.9, 1.65
    card_w, card_h = 10.5, 1.0
    for i, (title, content, role) in enumerate(layers):
        cy = y + i * (card_h + 0.15)
        add_box(slide, Inches(x), Inches(cy), Inches(card_w), Inches(card_h), BLUE_LIGHT if i % 2 == 0 else BLUE_PALE)
        add_textbox(slide, Inches(x + 0.2), Inches(cy + 0.12), Inches(1.6), Inches(0.35),
                    [[(title, 13, True, BLUE)]])
        add_textbox(slide, Inches(x + 2.0), Inches(cy + 0.18), Inches(5.5), Inches(0.3),
                    [[(content, 11, False, INK)]])
        add_box(slide, Inches(x + 8.0), Inches(cy + 0.2), Inches(2.2), Inches(0.5), BLUE if i == len(layers) - 1 else BLUE_MID)
        add_textbox(slide, Inches(x + 8.0), Inches(cy + 0.2), Inches(2.2), Inches(0.5),
                    [[(role, 10, True, WHITE)]], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(layers) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x + card_w / 2 - 0.15), Inches(cy + card_h + 0.02), Inches(0.3), Inches(0.14))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = BLUE_MID
            arrow.line.fill.background()
    # 核心原则
    principles = [
        "所有填报信息服务器持久化，换设备不丢失",
        "会议数据按单条 per-meeting 同步，避免误清",
        "向量库支持语义检索与历史知识复用",
        "关键操作保留版本审计，支持回溯",
    ]
    for i, p in enumerate(principles):
        py = 6.1 + (i // 2) * 0.45
        px = 0.9 + (i % 2) * 6.0
        add_textbox(slide, Inches(px), Inches(py), Inches(5.8), Inches(0.35),
                    [[("▪ ", 10, True, BLUE), (p, 10, False, MUTED)]])
    add_footer(slide)


def slide_ai_capabilities(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, "已落地的五大 AI 能力")
    caps = [
        ("文本理解", "纪要生成、决议提取、语义搜索、材料审核", "LLM（Kimi）"),
        ("推荐引擎", "AI 议程推荐、相关议题推荐、责任人推荐", "规则 + 协同过滤"),
        ("知识检索", "历史决议查询、相似案例匹配、KMS 问答", "RAG + 向量数据库"),
        ("智能预警", "指标偏离、逾期风险、会前准备度不足", "时序分析 + 阈值模型"),
        ("Agent 执行", "自动创建会议、分发行动项、规则触发通知", "Workflow + Tool Use"),
    ]
    # 表格
    rows = len(caps)
    x, y = 0.7, 1.65
    w, h = 12.0, 4.8
    col_ws = [2.2, 6.8, 3.0]
    # 表头
    headers = ["能力", "已上线场景", "技术形态"]
    header_h = 0.45
    add_box(slide, Inches(x), Inches(y), Inches(w), Inches(header_h), BLUE)
    cx = x
    for hi, hw in enumerate(col_ws):
        add_textbox(slide, Inches(cx + 0.1), Inches(y), Inches(hw - 0.2), Inches(header_h),
                    [[(headers[hi], 11, True, WHITE)]], valign=MSO_ANCHOR.MIDDLE)
        cx += hw
    # 行
    row_h = h / rows
    for ri, (cap, scenes, tech) in enumerate(caps):
        ry = y + header_h + ri * row_h
        fill = BLUE_LIGHT if ri % 2 == 0 else WHITE
        add_box(slide, Inches(x), Inches(ry), Inches(w), Inches(row_h - 0.03), fill)
        cx = x
        for ci, text in enumerate([cap, scenes, tech]):
            add_textbox(slide, Inches(cx + 0.1), Inches(ry + 0.08), Inches(col_ws[ci] - 0.2), Inches(row_h - 0.18),
                        [[(text, 10.5, ci == 0, INK if ci == 0 else TEXT)]], line_spacing=1.1)
            cx += col_ws[ci]
    # 底部统一底座标签
    add_box(slide, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.45), DARK_BG)
    add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11.8), Inches(0.45),
                [[("统一 AI 网关  ·  统一 AI 对话 UI  ·  全局 AI 抽屉「DSTE 智脑」", 10, False, WHITE)]],
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)


def slide_before_after(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, "传统模式 vs DSTE + AI", subtitle="从「人推着系统跑」到「系统推着人跑」")
    comparisons = [
        ("会前准备", "人工收集议题、整理材料、确认参会人，2–3 天", "AI 自动推荐议程、材料审核评分、参会人推荐，数小时"),
        ("会中记录", "专人记录，纪要会后 1–2 天产出", "会议 AI 助手实时辅助，关键信息自动提取"),
        ("决议跟进", "人工分发待办，逾期靠催办", "行动项自动分发、逾期提醒、progressLogs 跟进记录"),
        ("风险识别", "事后复盘才发现指标偏离", "规则引擎主动预警，会前/会中实时推送"),
        ("知识复用", "历史决议散落各处，检索困难", "KMS 知识库 + 语义检索，问答式调取"),
        ("数据口径", "多系统数据，会中对数耗时", "统一数据底座，会中随时问答取证"),
    ]
    x, y = 0.6, 1.55
    w = 12.1
    row_h = 0.68
    header_h = 0.42
    # 表头
    add_box(slide, Inches(x), Inches(y), Inches(w), Inches(header_h), BLUE)
    add_textbox(slide, Inches(x + 0.1), Inches(y), Inches(2.0), Inches(header_h),
                [[("维度", 11, True, WHITE)]], valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(x + 2.2), Inches(y), Inches(4.6), Inches(header_h),
                [[("传统模式", 11, True, WHITE)]], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(x + 6.9), Inches(y), Inches(5.0), Inches(header_h),
                [[("DSTE + AI（已上线）", 11, True, WHITE)]], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    # 行
    for i, (dim, before, after) in enumerate(comparisons):
        ry = y + header_h + i * row_h
        fill = BLUE_PALE if i % 2 == 0 else WHITE
        add_box(slide, Inches(x), Inches(ry), Inches(w), Inches(row_h - 0.03), fill)
        add_textbox(slide, Inches(x + 0.1), Inches(ry + 0.08), Inches(1.9), Inches(row_h - 0.16),
                    [[(dim, 10.5, True, BLUE)]], valign=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(x + 2.2), Inches(ry + 0.08), Inches(4.5), Inches(row_h - 0.16),
                    [[(before, 10, False, MUTED)]], line_spacing=1.05)
        add_box(slide, Inches(x + 6.85), Inches(ry + 0.06), Inches(5.1), Inches(row_h - 0.15), BLUE_LIGHT)
        add_textbox(slide, Inches(x + 6.95), Inches(ry + 0.1), Inches(4.9), Inches(row_h - 0.2),
                    [[(after, 10, False, INK)]], line_spacing=1.05)
    # 核心结论
    add_box(slide, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.55), BLUE)
    add_textbox(slide, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.55),
                [[("核心结论：人负责判断与决策，AI 负责收集、整理、提醒与推荐。", 12, True, WHITE)]],
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)


def slide_process(prs, title, subtitle, steps, highlight_idx=None):
    """通用流程图页（会前/会中/会后/日常）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, title, subtitle)
    n = len(steps)
    margin_x, gap = 0.55, 0.25
    card_w = (13.333 - 2 * margin_x - (n - 1) * gap) / n
    card_h = 4.6
    card_y = 1.85
    for i, (name, desc) in enumerate(steps):
        cx = margin_x + i * (card_w + gap)
        is_hl = highlight_idx is not None and i == highlight_idx
        fill = BLUE if is_hl else BLUE_LIGHT
        text = WHITE if is_hl else INK
        add_box(slide, Inches(cx), Inches(card_y), Inches(card_w), Inches(card_h), fill)
        # 序号圆
        num_c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + card_w / 2 - 0.18), Inches(card_y + 0.18), Inches(0.36), Inches(0.36))
        num_c.fill.solid()
        num_c.fill.fore_color.rgb = WHITE if is_hl else BLUE
        num_c.line.fill.background()
        add_textbox(slide, Inches(cx + card_w / 2 - 0.18), Inches(card_y + 0.18), Inches(0.36), Inches(0.36),
                    [[(str(i + 1), 12, True, BLUE if not is_hl else BLUE)]], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # 标题
        add_textbox(slide, Inches(cx + 0.12), Inches(card_y + 0.7), Inches(card_w - 0.24), Inches(0.45),
                    [[(name, 12.5, True, text)]], align=PP_ALIGN.CENTER)
        # 描述
        add_textbox(slide, Inches(cx + 0.12), Inches(card_y + 1.25), Inches(card_w - 0.24), Inches(card_h - 1.45),
                    [[(desc, 10, False, WHITE if is_hl else MUTED)]], line_spacing=1.1)
        # 箭头
        if i < n - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           Inches(cx + card_w + 0.03), Inches(card_y + card_h / 2 - 0.15),
                                           Inches(gap - 0.06), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = BLUE_MID
            arrow.line.fill.background()
    add_footer(slide)


def slide_daily_insights(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, "日常：从「开会看数」到「持续洞察」")
    items = [
        ("经营风险预警", "指标偏离目标时主动推送预警", BLUE),
        ("智能问答助手", "随时询问经营数据、历史决议、执行情况", BLUE),
        ("知识复用", "相似问题自动推荐历史会议与最佳实践", BLUE_MID),
        ("会议质量分析", "统计议题覆盖度、决议闭环率、行动项逾期率", BLUE_MID),
        ("全局会议待办面板", "跨会议聚合 7 类待办，统一入口推进", GREEN),
    ]
    # 仪表盘式三栏
    # 左：预警  中：待办  右：问答
    panels = [
        ("风险预警", "指标偏离 · 逾期风险 · 准备度不足", BLUE),
        ("全局待办", "报告/评审/纪要/决议/行动项/评估/流程推进", GREEN),
        ("智能问答", "经营数据 · 历史决议 · 执行情况 · 最佳实践", BLUE_MID),
    ]
    margin_x, gap = 0.6, 0.35
    col_w = (13.333 - 2 * margin_x - 2 * gap) / 3
    col_h = 4.6
    col_y = 1.8
    for i, (p_title, p_desc, p_color) in enumerate(panels):
        cx = margin_x + i * (col_w + gap)
        add_box(slide, Inches(cx), Inches(col_y), Inches(col_w), Inches(col_h), BLUE_LIGHT)
        # 顶部色条
        add_box(slide, Inches(cx), Inches(col_y), Inches(col_w), Inches(0.08), p_color, MSO_SHAPE.RECTANGLE)
        add_textbox(slide, Inches(cx + 0.15), Inches(col_y + 0.2), Inches(col_w - 0.3), Inches(0.45),
                    [[(p_title, 14, True, p_color)]], align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(cx + 0.15), Inches(col_y + 0.75), Inches(col_w - 0.3), Inches(col_h - 1.0),
                    [[(p_desc, 11, False, INK)]], align=PP_ALIGN.CENTER, line_spacing=1.15)
    add_footer(slide)


def slide_v7_highlights(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, "近期重要上新", subtitle="规则引擎 + 工作区 + 待办 + 目录管理")
    highlights = [
        ("规则引擎中心", "基于 KPI/指标定义落后/达成规则，手动执行、预览落后战区、生成会议草案", "从「人找问题」到「规则找人」"),
        ("驾驶舱多标签工作区", "支持跨阶段并行打开多个页面，标签状态持久化，iframe 嵌入外部页", "减少跳转，提升并行效率"),
        ("全局会议待办面板", "聚合报告/评审/纪要/决议/行动项/评估/流程推进 7 类待办", "会议工作一站式推进"),
        ("统一 Kimi 风格 AI 对话 UI", "全局 AI 抽屉、会议 AI 助手、专题 AI 统一气泡、快捷问题、流式输出", "体验一致，降低学习成本"),
        ("目录管理配置", "树形目录增删改、拖拽排序、启用/禁用，首期用于重点工作分类", "支持业务自定义分类体系"),
    ]
    # 2列卡片
    cols = 2
    margin_x, margin_y = 0.6, 1.65
    gap_x, gap_y = 0.35, 0.28
    card_w = (13.333 - 2 * margin_x - (cols - 1) * gap_x) / cols
    card_h = 0.95
    for i, (title, desc, value) in enumerate(highlights):
        c, r = i % cols, i // cols
        cx = margin_x + c * (card_w + gap_x)
        cy = margin_y + r * (card_h + gap_y)
        add_box(slide, Inches(cx), Inches(cy), Inches(card_w), Inches(card_h), BLUE_LIGHT)
        # 左侧色条
        add_box(slide, Inches(cx), Inches(cy), Inches(0.06), Inches(card_h), BLUE, MSO_SHAPE.RECTANGLE)
        add_textbox(slide, Inches(cx + 0.15), Inches(cy + 0.1), Inches(card_w - 0.3), Inches(0.3),
                    [[(title, 12.5, True, INK)]])
        add_textbox(slide, Inches(cx + 0.15), Inches(cy + 0.42), Inches(card_w - 0.3), Inches(0.28),
                    [[(desc, 9.5, False, MUTED)]], line_spacing=1.05)
        # 价值标签
        add_box(slide, Inches(cx + card_w - 2.6), Inches(cy + 0.55), Inches(2.5), Inches(0.32), BLUE_PALE)
        add_textbox(slide, Inches(cx + card_w - 2.6), Inches(cy + 0.55), Inches(2.5), Inches(0.32),
                    [[(value, 8.5, False, BLUE)]], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, "当前建设成果")
    metrics = [
        ("产品版本", "v0.7.14"),
        ("E2E 测试覆盖", "359 passed / 25 skipped"),
        ("核心模块", "战略地图、RoadMap、年度计划、OMP、经营分析会、决议中心、人员组织、规则引擎、驾驶舱工作区"),
        ("AI 落地场景", "议程推荐、材料审核、纪要生成、决议提取、行动项分发、KMS 问答、全局 AI 抽屉、智能预警"),
        ("数据同步机制", "业务/战略/OMP/洞察/议题/评审已接入 per-record 云端同步"),
        ("关键效率提升", "会前准备 2–3 天 → 数小时；会议纪要 1–2 天 → 5 分钟"),
    ]
    # 上半区：大数字卡片
    big_cards = [
        ("v0.7.14", "当前版本", BLUE),
        ("359", "E2E 用例通过", GREEN),
        ("8+", "AI 场景落地", ORANGE),
    ]
    margin_x, gap = 0.8, 0.4
    card_w = (13.333 - 2 * margin_x - 2 * gap) / 3
    for i, (num, label, color) in enumerate(big_cards):
        cx = margin_x + i * (card_w + gap)
        add_box(slide, Inches(cx), Inches(1.75), Inches(card_w), Inches(1.6), BLUE_LIGHT)
        add_textbox(slide, Inches(cx), Inches(1.9), Inches(card_w), Inches(0.9),
                    [[(num, 32, True, color)]], align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(cx), Inches(2.75), Inches(card_w), Inches(0.4),
                    [[(label, 11, False, MUTED)]], align=PP_ALIGN.CENTER)
    # 下半区：详细指标表
    x, y = 0.7, 3.65
    w, h = 12.0, 3.0
    row_h = h / len(metrics)
    for i, (label, value) in enumerate(metrics):
        ry = y + i * row_h
        fill = BLUE_PALE if i % 2 == 0 else WHITE
        add_box(slide, Inches(x), Inches(ry), Inches(w), Inches(row_h - 0.03), fill)
        add_textbox(slide, Inches(x + 0.15), Inches(ry + 0.06), Inches(2.5), Inches(row_h - 0.12),
                    [[(label, 10.5, True, BLUE)]], valign=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(x + 2.75), Inches(ry + 0.06), Inches(w - 2.9), Inches(row_h - 0.12),
                    [[(value, 10.5, False, INK)]], valign=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    # 总结条
    add_box(slide, Inches(0.7), Inches(6.75), Inches(12.0), Inches(0.42), BLUE_LIGHT)
    add_textbox(slide, Inches(0.8), Inches(6.75), Inches(11.8), Inches(0.42),
                [[("核心模块已齐套，进入打磨与规模化阶段", 11, True, INK)]],
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)


def slide_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, "下阶段重点：从「能用」到「好用」")
    # 左侧：近期重点
    add_box(slide, Inches(0.6), Inches(1.55), Inches(6.2), Inches(4.7), BLUE_PALE)
    add_textbox(slide, Inches(0.75), Inches(1.65), Inches(5.9), Inches(0.45),
                [[("近期重点（未来 1–2 个月）", 14, True, BLUE)]])
    steps = [
        "AI 效果可量化：建立使用数据看板，追踪采纳率、准确率、效率提升",
        "场景深度打磨：聚焦经营分析会全链路，优化议程推荐、纪要质量、决议闭环",
        "规模化推广：确认首批试点业务线，收集真实反馈，沉淀最佳实践",
        "稳定性与性能：持续补齐 E2E 覆盖，优化大文件/大数据量场景性能",
    ]
    for i, s in enumerate(steps):
        add_textbox(slide, Inches(0.85), Inches(2.2 + i * 0.85), Inches(5.8), Inches(0.75),
                    [[(f"{i+1}. ", 11, True, BLUE), (s, 10.5, False, INK)]], line_spacing=1.1)
    # 右侧：风险与应对
    add_box(slide, Inches(7.1), Inches(1.55), Inches(5.6), Inches(4.7), BLUE_LIGHT)
    add_textbox(slide, Inches(7.25), Inches(1.65), Inches(5.3), Inches(0.45),
                [[("风险与应对", 14, True, BLUE)]])
    risks = [
        ("AI 推荐采纳率不及预期", "建立人工反馈闭环，A/B 测试推荐策略"),
        ("数据质量影响 AI 效果", "强化数据治理与录入校验，推进源系统对接"),
        ("多模块并行维护成本高", "持续重构公共组件与设计系统，减少重复代码"),
    ]
    for i, (risk, deal) in enumerate(risks):
        ry = 2.25 + i * 1.1
        add_textbox(slide, Inches(7.35), Inches(ry), Inches(5.2), Inches(0.35),
                    [[("风险：", 10, True, ORANGE), (risk, 10, False, INK)]], line_spacing=1.05)
        add_textbox(slide, Inches(7.35), Inches(ry + 0.35), Inches(5.2), Inches(0.4),
                    [[("应对：", 10, True, GREEN), (deal, 10, False, MUTED)]], line_spacing=1.05)
    # 下阶段目标
    add_box(slide, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.55), BLUE)
    add_textbox(slide, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.55),
                [[("下阶段目标：体验好用、效果可量化、场景可规模化", 12, True, WHITE)]],
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, MSO_SHAPE.RECTANGLE)
    # 右侧装饰
    add_box(slide, Inches(8.8), Inches(0), Inches(4.53), prs.slide_height, BLUE_PALE, MSO_SHAPE.RECTANGLE)
    add_section_title(slide, "DSTE + AI：让每一次战略执行都更有价值")
    conclusions = [
        "DSTE 已从单一工具升级为覆盖「战略 → 执行 → 分析 → 洞察」的战略执行平台",
        "AI 能力已贯通会前、会中、会后与日常经营，进入规模化落地阶段",
        "数据底座与知识库是 AI 效果的基础，per-record 同步机制已基本补齐",
        "下阶段核心是从「功能可用」迈向「体验好用、效果可量化」",
    ]
    for i, c in enumerate(conclusions):
        add_textbox(slide, Inches(0.75), Inches(2.0 + i * 0.75), Inches(7.5), Inches(0.6),
                    [[(f"{i+1}. ", 12, True, BLUE), (c, 11.5, False, INK)]], line_spacing=1.1)
    # 右侧 checklist
    add_textbox(slide, Inches(9.0), Inches(1.8), Inches(4.0), Inches(0.4),
                [[("下一步行动", 14, True, BLUE)]])
    actions = [
        "确认首批试点业务线与核心场景",
        "建立 AI 功能使用数据看板与效果评估机制",
        "持续打磨经营分析会全流程体验",
        "推进数据治理与源系统对接",
    ]
    for i, a in enumerate(actions):
        add_box(slide, Inches(9.0), Inches(2.35 + i * 0.65), Inches(0.22), Inches(0.22), BLUE, MSO_SHAPE.RECTANGLE)
        add_textbox(slide, Inches(9.35), Inches(2.32 + i * 0.65), Inches(3.8), Inches(0.5),
                    [[(a, 10.5, False, INK)]], line_spacing=1.05)
    add_footer(slide)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_agenda(prs)
    slide_pain_points(prs)
    slide_value_positioning(prs)
    slide_milestones(prs)
    slide_architecture(prs)
    slide_data_architecture(prs)
    slide_ai_capabilities(prs)
    slide_before_after(prs)
    slide_process(prs, "会前：从「催办」到「智能推荐」", None,
                   steps=[
                       ("议题采集", "AI 自动从 RoadMap、OKR、上期决议中抓取待议事项"),
                       ("议程推荐", "基于重要度、紧急度、关联性排序，给出置信度与推荐理由"),
                       ("材料审核", "4 套场景化评分矩阵，输出总分、维度打分、问题清单、改进建议"),
                       ("参会人推荐", "根据议题关联的角色与部门，推荐必要参会人"),
                       ("准备度检查", "自动计算准备度清单，材料审核 ≥60 分才算通过"),
                   ])
    slide_process(prs, "会中：从「记录」到「实时智能辅助」", None,
                   steps=[
                       ("会议 AI 助手", "注入当前会议上下文，流式问答，附快捷问题芯片"),
                       ("工具调用", "实时查询议程、行动项、决议执行情况，边开会边取证"),
                       ("关键信息提取", "自动标出决议、行动项、责任人、Deadline"),
                       ("数据即时问答", "会中随时问“Q3 华东区营收达成率多少？”"),
                       ("历史关联推送", "自动关联往期相似决议与执行情况"),
                   ])
    slide_process(prs, "会后：从「纪要」到「闭环追踪」", None,
                   steps=[
                       ("纪要自动生成", "结构化输出：议题 → 讨论要点 → 决议 → 行动项"),
                       ("决议去重归档", "AI 识别重复/相似决议，自动合并或关联历史"),
                       ("行动项分发", "按责任人推送待办，接入企业微信/钉钉"),
                       ("进展智能同步", "执行人填报后自动汇总，progressLogs 记录跟进"),
                       ("闭环追踪问答", "一句话查询未闭环行动项、决议执行状态"),
                   ])
    slide_daily_insights(prs)
    slide_v7_highlights(prs)
    slide_results(prs)
    slide_next_steps(prs)
    slide_summary(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"已生成: {OUT}（共 {len(prs.slides)} 页）")


if __name__ == "__main__":
    build_deck()
