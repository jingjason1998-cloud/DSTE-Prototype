#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""把多个 .pptx 合并成一个（保留版式、文字、图片）。

用法:
    .pptx-venv/bin/python scripts/merge_pptx.py 输出.pptx 输入1.pptx 输入2.pptx ...

说明:
    - 以 16:9（12192000 x 6858000 EMU）为目标画布；
    - 逐页复制形状元素，并重映射图片关系（r:embed），跨包拷贝图片 part；
    - 纯文本/形状页面（无媒体）直接深拷贝即可。
"""
import copy
import io
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

EMU_W, EMU_H = 12192000, 6858000  # 16:9


def copy_slide(dest_prs, src_slide):
    blank_layout = dest_prs.slide_layouts[6]
    dest = dest_prs.slides.add_slide(blank_layout)
    # 清掉版式自带占位符
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    # 复制页面背景（如有）
    src_bg = src_slide._element.find(qn("p:cSld") + "/" + qn("p:bg"))
    if src_bg is None:
        src_csld = src_slide._element.find(qn("p:cSld"))
        if src_csld is not None:
            src_bg = src_csld.find(qn("p:bg"))
    if src_bg is not None:
        dest_csld = dest._element.find(qn("p:cSld"))
        dest_csld.insert(0, copy.deepcopy(src_bg))

    # 源版式中的图片（如整页背景图）衬到最底层
    layout_pics = [
        sp for sp in src_slide.slide_layout._element.find(qn("p:cSld") + "/" + qn("p:spTree"))
        if sp.tag == qn("p:pic")
    ] if src_slide.slide_layout._element.find(qn("p:cSld") + "/" + qn("p:spTree")) is not None else []
    for pic in layout_pics:
        dest.shapes._spTree.insert(2, copy.deepcopy(pic))

    # 复制全部形状
    for shp in src_slide.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))

    # 重映射图片关系：rId -> 新 rId（页面自身 + 版式衬底图）
    src_parts = [src_slide.part]
    if layout_pics:
        src_parts.append(src_slide.slide_layout.part)
    for src_part in src_parts:
        for src_rid, rel in src_part.rels.items():
            if "image" not in rel.reltype:
                continue
            img_blob = rel.target_part.blob
            _, new_rid = dest.part.get_or_add_image_part(io.BytesIO(img_blob))
            for blip in dest._element.iter(qn("a:blip")):
                if blip.get(qn("r:embed")) == src_rid:
                    blip.set(qn("r:embed"), new_rid)


def merge(out_path: Path, inputs: list):
    dest = Presentation()
    dest.slide_width = EMU_W
    dest.slide_height = EMU_H
    # 删除默认空白页（若有）
    while len(dest.slides) > 0:
        r_id = dest.slides._sldIdLst[0].get(qn("r:id"))
        dest.part.drop_rel(r_id)
        del dest.slides._sldIdLst[0]

    for f in inputs:
        src = Presentation(f)
        for slide in src.slides:
            copy_slide(dest, slide)
        print(f"  + {f}（{len(src.slides)} 页）")

    dest.save(out_path)
    print(f"已生成: {out_path}（共 {len(Presentation(str(out_path)).slides)} 页）")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)
    merge(Path(sys.argv[1]), sys.argv[2:])
