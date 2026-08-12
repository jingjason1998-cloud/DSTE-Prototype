#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""重做 DSTE汇报合集.pptx 前 4 页（pptx-design skill 版，含图标与语义色）。

每页一种版式，杜绝卡片阵列套路：
- P1 事不过三机制  → 阶梯递进：顺延 1→2→3 逐级升高加深，步内嵌 Phosphor 图标
- P2 会议智能评价  → 比例评分条：35/30/35 分段条 + 明细列对齐，等级用绿/蓝/橙/红语义色
- P3 场景与价值    → 竖向时间轴（图标节点）+ hairline 表格（彩色状态 pill）+ 图标四价值
- P4 AI 能力总览   → 编辑式三栏：无卡片，竖 hairline 分隔，阶段图标 + 幽灵数字

图标由 scripts/build_icon_pngs.mjs 生成（scripts/assets/icons/）。

用法:
    .pptx-venv/bin/python scripts/redesign_base_pages.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
DECK = REPO_ROOT / "DSTE汇报合集.pptx"
ICON_DIR = REPO_ROOT / "scripts/assets/icons"

FONT = "PingFang SC"

BLUE = RGBColor(0x03, 0x5D, 0xCF)
BLUE_MID = RGBColor(0x5B, 0x8C, 0xC8)
CARD_LIGHT = RGBColor(0xE8, 0xF4, 0xFD)
CARD_MID = RGBColor(0xBF, 0xDC, 0xFC)
GHOST = RGBColor(0xD8, 0xE9, 0xFB)
INK = RGBColor(0x26, 0x26, 0x26)
INK_SOFT = RGBColor(0x3D, 0x4A, 0x5C)
MUTED = RGBColor(0x59, 0x59, 0x59)
HAIRLINE = RGBColor(0xD9, 0xDE, 0xE4)
WARN = RGBColor(0xE0, 0x52, 0x52)
WARN_BG = RGBColor(0xFD, 0xEC, 0xEA)
GREEN = RGBColor(0x2B, 0xA4, 0x71)
ORANGE = RGBColor(0xF2, 0x99, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

MARGIN = 0.72
PAGE_W = 13.333
CONTENT_W = PAGE_W - 2 * MARGIN  # 11.893
BRAND = "DSTE 战略管理平台 · 经分会智能闭环"


# ---------------------------------------------------------------- 基础工具
def set_text(tf, paras, align=PP_ALIGN.LEFT, space_after=Pt(4), line_spacing=1.0, spacing=None):
    tf.word_wrap = True
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.line_spacing = line_spacing
        for text, size, bold, color in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            if spacing:
                r.font._rPr.set("spc", str(spacing))


def add_shape(slide, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=None, **kw):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor:
        tf.vertical_anchor = anchor
    set_text(tf, paras, align=align, **kw)
    return box


def add_icon(slide, name, color, x, y, size):
    path = ICON_DIR / f"{name}-{color}.png"
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(size), Inches(size))


def add_pill(slide, x, y, w, h, fill, text, size, color, line=None, line_w=1.0):
    """文字直接写进形状的 text_frame 并垂直居中，保证与背景框协调。"""
    sp = add_shape(slide, x, y, w, h, fill, MSO_SHAPE.ROUNDED_RECTANGLE, line=line, line_w=line_w)
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(tf, [[(text, size, True, color)]], align=PP_ALIGN.CENTER)
    return sp


def header(slide, eyebrow, title_runs, subtitle):
    add_text(slide, MARGIN, 0.5, 10, 0.3, [[(eyebrow, 11, False, MUTED)]], spacing=200)
    add_text(slide, MARGIN, 0.88, 11.9, 0.65, [title_runs])
    add_text(slide, MARGIN, 1.53, 11.9, 0.35, [[(subtitle, 13, False, MUTED)]])


def footer(slide, right_text="", page=""):
    add_shape(slide, MARGIN, 6.98, CONTENT_W, 0.012, HAIRLINE)
    add_text(slide, MARGIN, 7.08, 6, 0.3, [[(BRAND, 9, False, MUTED)]])
    tail = f"{right_text}　{page}" if right_text else page
    if tail:
        add_text(slide, PAGE_W - MARGIN - 6, 7.08, 6, 0.3,
                 [[(tail, 9.5, False, MUTED)]], align=PP_ALIGN.RIGHT)


def label_with_icon(slide, x, y, icon, text, w=3.6):
    add_icon(slide, icon, "035DCF", x, y + 0.015, 0.2)
    add_text(slide, x + 0.3, y, w, 0.3, [[(text, 11, True, BLUE)]])


def blank_layout(prs):
    layouts = prs.slide_masters[0].slide_layouts
    return next((l for l in layouts if l.name in ("空白", "Blank")), layouts[6])


def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    sldId = list(xml_slides)[index]
    rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rId)
    xml_slides.remove(sldId)


# ------------------------------------------------- P1 事不过三机制：阶梯递进
def build_p1(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header(slide, "经分会 · 议程管理",
           [("「事不过三」", 28, True, BLUE), ("议程顺延机制", 28, True, INK)],
           "把经营分析会，变成干部的试金石")

    steps = [
        ("1", "第一次顺延", "记录原因与责任人", CARD_LIGHT, INK, MUTED, 1.15, "pencil-simple", "035DCF"),
        ("2", "第二次顺延", "累计历史，链路自动关联", CARD_MID, INK, INK_SOFT, 1.65, "link", "035DCF"),
        ("3", "第三次顺延", "强制警示，升级高管督办", BLUE, WHITE, CARD_LIGHT, 2.15, "siren", "FFFFFF"),
    ]
    step_w, gap, bottom = 3.74, 0.37, 4.75
    for i, (num, label, desc, fill, c_name, c_desc, h, icon, icon_color) in enumerate(steps):
        x = MARGIN + i * (step_w + gap)
        y = bottom - h
        add_shape(slide, x, y, step_w, h, fill, MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(slide, x + 0.26, y + 0.13, step_w - 0.9, 0.5,
                 [[(num + "  ", 24, True, c_name), (label, 14, True, c_name)]])
        add_icon(slide, icon, icon_color, x + step_w - 0.62, y + 0.16, 0.32)
        add_text(slide, x + 0.26, bottom - 0.42, step_w - 0.52, 0.32,
                 [[(desc, 9.5, False, c_desc)]])

    # 底部三栏：机制逻辑 / 产品能力 / 价值主张
    top = 5.2
    label_with_icon(slide, MARGIN, top, "gear", "机制逻辑")
    add_text(slide, MARGIN, top + 0.34, 3.6, 1.3,
             [[("议题第一次讲不清楚，允许顺延；连续三次仍无结论，系统自动标红警示 —— 倒逼干部会前想透、会中讲清、会后闭环。", 10, False, INK)]],
             line_spacing=1.3)

    label_with_icon(slide, 4.72, top, "list-checks", "产品能力")
    add_text(slide, 4.72, top + 0.34, 3.6, 1.3,
             [[("▪ ", 9.5, True, BLUE), ("议程状态标记：已完成 / 已顺延 / 已承接", 9.5, False, INK)],
              [("▪ ", 9.5, True, BLUE), ("顺延次数自动累计，历史链路可追溯", 9.5, False, INK)],
              [("▪ ", 9.5, True, BLUE), ("第三次顺延触发强制警示", 9.5, False, INK)],
              [("▪ ", 9.5, True, BLUE), ("来源会议自动关联，议题不空转", 9.5, False, INK)]],
             space_after=Pt(3), line_spacing=1.15)

    add_shape(slide, 8.72, top + 0.05, 0.03, 1.35, BLUE)
    label_with_icon(slide, 8.95, top, "lightbulb", "价值主张")
    add_text(slide, 8.95, top + 0.34, 3.68, 1.35,
             [[("一个问题讲三遍还讲不清、没有解决方案，说明干部没想透、没担当、没能力 —— 用数据把它钉在桌面上。", 10, False, INK)]],
             line_spacing=1.3)

    footer(slide, "「三遍讲不清，干部要负责」", "01 / 06")


# ------------------------------------------------- P2 智能评价：比例评分条
def build_p2(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header(slide, "经分会 · 会议评价",
           [("本部会议", 28, True, INK), ("智能评价", 28, True, BLUE)],
           "让会议质量可量化、可对比、可改进")

    # 顺延扣分警示标签（评分条右上方）
    add_pill(slide, 10.1, 1.78, 2.51, 0.32, WHITE, "任一议程顺延 −5 分", 9.5, WARN,
             line=WARN, line_w=1.0)

    # 35 / 30 / 35 比例分段条（段内嵌图标）
    bar_y, bar_h = 2.2, 0.8
    segs = [
        ("会前 · 35 分", 0.35, CARD_LIGHT, INK, "clipboard-text", "035DCF"),
        ("会中 · 30 分", 0.30, CARD_MID, INK, "chat", "035DCF"),
        ("会后 · 35 分", 0.35, BLUE, WHITE, "flag", "FFFFFF"),
    ]
    x = MARGIN
    for label, ratio, fill, c, icon, icon_color in segs:
        w = CONTENT_W * ratio
        add_shape(slide, x, bar_y, w, bar_h, fill)
        add_text(slide, x, bar_y, w, bar_h, [[(label, 14, True, c)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 图标放在居中文字的左侧（文字约 8 字符 ≈ 1.55in）
        add_icon(slide, icon, icon_color, x + w / 2 - 1.15, bar_y + 0.25, 0.3)
        x += w

    # 明细列与分段对齐
    details = [
        ["材料完整 7 分", "议程覆盖 10.5 分", "材料评分 17.5 分"],
        ["有效讨论 12 分", "参与度 12 分", "时间控制 6 分"],
        ["有效决议与待办 30 分", "评分及时性 5 分"],
    ]
    x = MARGIN
    for (label, ratio, fill, c, icon, icon_color), items in zip(segs, details):
        w = CONTENT_W * ratio
        add_text(slide, x + 0.15, bar_y + bar_h + 0.18, w - 0.3, 0.9,
                 [[(it, 9.5, False, MUTED)] for it in items],
                 align=PP_ALIGN.CENTER, space_after=Pt(3), line_spacing=1.15)
        x += w

    # 等级标准图例（语义色：绿 / 蓝 / 橙 / 红）
    grades = [("优秀 ≥90", GREEN), ("良好 ≥75", BLUE), ("及格 ≥60", ORANGE), ("待改进 <60", WARN)]
    add_text(slide, MARGIN, 4.32, 1.2, 0.3, [[("等级标准", 10, True, INK)]])
    gx = 1.9
    for name, c in grades:
        add_shape(slide, gx, 4.36, 0.16, 0.16, c, MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(slide, gx + 0.24, 4.30, 1.6, 0.3, [[(name, 10, False, INK)]])
        gx += 1.85

    # 底部两栏：价值主张 / AI 的作用
    top = 4.95
    add_shape(slide, 6.6, top + 0.05, 0.012, 1.65, HAIRLINE)
    label_with_icon(slide, MARGIN, top, "lightbulb", "价值主张", w=5.5)
    add_text(slide, MARGIN, top + 0.34, 5.5, 1.5,
             [[("▪ ", 9.5, True, BLUE), ("把会议从「开完就算」变成「可评可改」", 9.5, False, INK)],
              [("▪ ", 9.5, True, BLUE), ("量化会前准备、会中执行、会后闭环", 9.5, False, INK)],
              [("▪ ", 9.5, True, BLUE), ("横向对比各本部，识别优秀与落后", 9.5, False, INK)],
              [("▪ ", 9.5, True, BLUE), ("为干部考核与会议改进提供数据依据", 9.5, False, INK)]],
             space_after=Pt(3), line_spacing=1.15)
    label_with_icon(slide, 6.95, top, "robot", "AI 的作用", w=5.6)
    add_text(slide, 6.95, top + 0.34, 5.6, 1.5,
             [[("▪ ", 9.5, True, BLUE), ("自动采集会议数据，实时计算三阶段得分", 9.5, False, INK)],
              [("▪ ", 9.5, True, BLUE), ("规则生成反馈标签：材料充分 / 讨论有效 / 闭环到位", 9.5, False, INK)],
              [("▪ ", 9.5, True, BLUE), ("识别低分会议与风险议题，推送改进建议", 9.5, False, INK)],
              [("▪ ", 9.5, True, BLUE), ("减少人工打分偏差，让评价客观可复现", 9.5, False, INK)]],
             space_after=Pt(3), line_spacing=1.15)

    footer(slide, "「会议好不好，数据说了算」", "02 / 06")


# ------------------------------------------- P3 场景与价值：时间轴 + 表格
def build_p3(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header(slide, "经分会 · 场景与价值",
           [("让「讲不清」的议题", 26, True, BLUE), ("，在数据看板上无处藏身", 26, True, INK)],
           "场景示例：华东区 Q2 营收下滑分析的三次经分会")

    # 左：竖向时间轴（图标节点）
    line_x = 1.19
    add_shape(slide, line_x, 2.5, 0.016, 2.3, HAIRLINE)
    timeline = [
        ("第一次经分会", "已顺延", "罗列了大量数据，但未找到根因与责任人", False),
        ("第二次经分会", "已顺延", "归因市场与竞品，仍无解决方案与行动计划", False),
        ("第三次经分会", "已闭环", "系统标红警示，高管督办，形成闭环行动", True),
    ]
    for i, (name, status, desc, done) in enumerate(timeline):
        y = 2.3 + i * 1.05
        node_fill = GREEN if done else BLUE
        add_shape(slide, line_x - 0.163, y, 0.34, 0.34, node_fill, MSO_SHAPE.OVAL)
        add_icon(slide, "calendar-check", "FFFFFF", line_x - 0.093, y + 0.07, 0.2)
        add_text(slide, 1.55, y, 1.6, 0.32, [[(name, 12, True, INK)]])
        if done:
            add_pill(slide, 3.2, y + 0.01, 0.92, 0.3, GREEN, status, 9, WHITE)
        else:
            add_pill(slide, 3.2, y + 0.01, 0.92, 0.3, CARD_LIGHT, status, 9, BLUE)
        add_text(slide, 1.55, y + 0.38, 3.3, 0.55, [[(desc, 9.5, False, MUTED)]], line_spacing=1.2)

    # 右：hairline 数据表（彩色状态 pill）
    tx, tw = 5.15, 7.46
    cols = [(0.0, 3.2, "议题"), (3.3, 1.05, "责任人"), (4.45, 1.1, "顺延次数"), (5.65, 1.81, "当前状态")]
    add_text(slide, tx, 2.02, 4, 0.3, [[("数据看板 · 议题闭环跟踪", 11, True, BLUE)]])
    for cx, cw, name in cols:
        add_text(slide, tx + cx, 2.42, cw, 0.28, [[(name, 9.5, True, MUTED)]])
    add_shape(slide, tx, 2.74, tw, 0.012, HAIRLINE)

    rows = [
        ("华东区 Q2 营收下滑根因分析", "张经理", "3", "⚠ 已警示", WARN_BG, WARN, None),
        ("供应链降本专项推进", "李总监", "1", "进行中", None, BLUE, BLUE),
        ("新品上市节奏复盘", "王总监", "0", "已完成", GREEN, WHITE, None),
        ("客户满意度下滑应对", "赵经理", "2", "已顺延", CARD_LIGHT, BLUE, None),
    ]
    for i, (topic, owner, cnt, status, s_fill, s_text, s_line) in enumerate(rows):
        y = 2.88 + i * 0.58
        add_text(slide, tx, y, 3.2, 0.3, [[(topic, 10.5, False, INK)]])
        add_text(slide, tx + 3.3, y, 1.05, 0.3, [[(owner, 10.5, False, INK)]])
        add_text(slide, tx + 4.45, y, 1.1, 0.3, [[(cnt, 10.5, True, BLUE)]])
        add_pill(slide, tx + 5.65, y - 0.02, 1.05, 0.32, s_fill, status, 9, s_text,
                 line=s_line, line_w=1.0)
        if i < 3:
            add_shape(slide, tx, y + 0.44, tw, 0.008, HAIRLINE)

    # 底部：图标四价值（竖 hairline 分隔）
    add_text(slide, MARGIN, 5.35, 6, 0.3, [[("产品功能带来的价值", 11, True, BLUE)]])
    values = [
        ("gauge", "提升会议效率", "减少无效议题反复占用高管时间"),
        ("shield-check", "强化干部担当", "用数据暴露「讲不清」的问题"),
        ("link", "沉淀决策链路", "议题历史可追溯，避免推诿空转"),
        ("check-circle", "加速问题闭环", "倒逼会前准备，缩短决策周期"),
    ]
    vw = CONTENT_W / 4
    for i, (icon, name, desc) in enumerate(values):
        x = MARGIN + i * vw
        if i > 0:
            add_shape(slide, x, 5.78, 0.012, 0.95, HAIRLINE)
        ix = x + (0.25 if i > 0 else 0)
        add_icon(slide, icon, "035DCF", ix, 5.74, 0.26)
        add_text(slide, ix + 0.36, 5.74, vw - 0.65, 0.3, [[(name, 12, True, INK)]])
        add_text(slide, ix, 6.14, vw - 0.4, 0.55, [[(desc, 9, False, MUTED)]], line_spacing=1.2)

    footer(slide, "", "03 / 06")


# ------------------------------------------------- P4 AI 总览：编辑式三栏
def build_p4(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header(slide, "经营分析会 · AI 能力",
           [("经营分析会全流程 ", 28, True, INK), ("AI 能力", 28, True, BLUE)],
           "会前准备 · 会中陪伴 · 会后闭环，统一 AI 底座支撑")

    phases = [
        ("01", "会前", "议程与材料准备", "calendar-check", [
            ("AI 议程推荐", "融合历史议程、顺延议题与未闭环行动项，给出置信度与推荐理由，一键写入议程"),
            ("材料智能审核（KMS）", "4 套场景化评分矩阵，输出总分、打分理由、问题清单、改进建议与审核结论"),
            ("一键 / 批量送审", "议程行直接送审，评分自动回传彩色徽标，会议详情页汇总平均分"),
            ("会前准备度检查", "自动计算准备度清单，材料审核评分 ≥60 才算通过"),
        ]),
        ("02", "会中", "智能陪伴与执行", "robot", [
            ("会议 AI 助手", "注入当前会议上下文的流式问答，附快捷问题芯片，断网降级本地规则回复"),
            ("工具调用", "实时查询议程、行动项、决议执行情况，边开会边取证"),
            ("AI 草拟 · 人工确认", "AI 生成行动项 / 新会议草案卡片，点击确认后才写入，可控可追溯"),
        ]),
        ("03", "会后", "评估与闭环追踪", "arrows-clockwise", [
            ("AI 自动评分", "三段式模型（35+30+35），直接消费材料审核分，议程顺延自动扣分"),
            ("闭环追踪问答", "一句话查询未闭环行动项与决议状态，自动生成经营分析周报草稿"),
        ]),
    ]
    col_w = 3.75
    col_xs = [MARGIN, 4.85, 8.98]
    for lx in (4.55, 8.68):
        add_shape(slide, lx, 2.25, 0.012, 4.35, HAIRLINE)

    for (num, name, tag, icon, items), cx in zip(phases, col_xs):
        add_icon(slide, icon, "035DCF", cx, 2.24, 0.34)
        add_text(slide, cx + 0.5, 2.2, col_w - 1.5, 0.45,
                 [[(name, 20, True, INK), ("  " + tag, 11, False, MUTED)]])
        add_text(slide, cx + col_w - 1.0, 2.06, 1.0, 0.6,
                 [[(num, 34, True, GHOST)]], align=PP_ALIGN.RIGHT)
        add_shape(slide, cx, 2.78, col_w, 0.014, HAIRLINE)
        paras = []
        for iname, idesc in items:
            paras.append([("▪ ", 10, True, BLUE), (iname, 11, True, INK)])
            paras.append([(idesc, 9.5, False, MUTED)])
        add_text(slide, cx, 2.98, col_w, 3.7, paras, space_after=Pt(6), line_spacing=1.2)

    add_shape(slide, MARGIN, 6.98, CONTENT_W, 0.012, HAIRLINE)
    add_text(slide, MARGIN, 7.08, 10.5, 0.3,
             [[("AI 底座 ｜ ", 9.5, True, BLUE),
               ("统一 AI 网关（Kimi 大模型，SSE 流式） · KMS 知识库工具 · 全局 AI 抽屉「DSTE 智脑」 · 统一 AI 交互 UI", 9.5, False, MUTED)]])
    add_text(slide, PAGE_W - MARGIN - 1.2, 7.08, 1.2, 0.3,
             [[("04 / 06", 9.5, False, MUTED)]], align=PP_ALIGN.RIGHT)


def main():
    prs = Presentation(str(DECK))
    build_p1(prs)
    build_p2(prs)
    build_p3(prs)
    build_p4(prs)
    # 新 4 页移到最前，再删掉旧 4 页
    xml = prs.slides._sldIdLst
    ids = list(xml)
    new_ids = ids[-4:]
    for sld in new_ids:
        xml.remove(sld)
    for i, sld in enumerate(new_ids):
        xml.insert(i, sld)
    for _ in range(4):
        delete_slide(prs, 4)
    prs.save(str(DECK))
    print(f"已重做前 4 页: {DECK}（共 {len(prs.slides._sldIdLst)} 页）")


if __name__ == "__main__":
    main()
