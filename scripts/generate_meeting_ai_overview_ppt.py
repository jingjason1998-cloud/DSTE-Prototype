#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成《经营分析会管理全流程 AI 功能总览》单页 PPT。

用法:
    .pptx-venv/bin/python scripts/generate_meeting_ai_overview_ppt.py [输出路径]

默认输出: 仓库根目录 经营分析会AI功能总览.pptx
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
DEFAULT_OUT = REPO_ROOT / "经营分析会AI功能总览.pptx"

FONT = "PingFang SC"

# 配色
INK = RGBColor(0x1F, 0x29, 0x37)       # 主文字
MUTED = RGBColor(0x6B, 0x72, 0x80)     # 次要文字
BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xF7, 0xF8, 0xFA)
BASE_BG = RGBColor(0x16, 0x1B, 0x26)   # 底座深色

PHASES = [
    {
        "name": "会前",
        "sub": "议程与材料准备",
        "color": RGBColor(0x1F, 0x6F, 0xEB),
        "items": [
            ("AI 议程推荐", "融合历史议程、顺延议题、未闭环行动项/决议与 OMP 重点工作，给出置信度与推荐理由，人工采纳后一键写入议程"),
            ("材料智能审核（KMS）", "4 套场景化评分矩阵，输出总分、维度打分理由、问题清单、改进建议、亮点与审核结论，支持批量审核与历史对比"),
            ("一键 / 批量送审", "议程行直接送审，评分自动回传并显示彩色得分徽标，会议详情页汇总平均分"),
            ("会前准备度检查", "自动计算准备度清单，材料审核评分 ≥60 才算通过"),
        ],
    },
    {
        "name": "会中",
        "sub": "智能陪伴与执行",
        "color": RGBColor(0x7C, 0x3A, 0xED),
        "items": [
            ("会议 AI 助手", "注入当前会议上下文（议程/纪要/行动项/决议）的流式问答，附快捷问题芯片，断网降级为本地规则回复"),
            ("工具调用", "实时查询议程、行动项、决议执行情况，边开会边取证"),
            ("AI 草拟 · 人工确认", "AI 生成行动项/新会议草案卡片，点击确认后才写入，可控可追溯"),
        ],
    },
    {
        "name": "会后",
        "sub": "评估与闭环追踪",
        "color": RGBColor(0x0E, 0x9F, 0x6E),
        "items": [
            ("AI 自动评分", "三段式模型（会前 35 + 会中 30 + 会后 35），直接消费材料审核分，议程顺延自动扣分，生成反馈标签"),
            ("闭环追踪问答", "一句话查询未闭环行动项、决议执行状态，自动生成经营分析周报草稿"),
        ],
    },
]

BASE_FEATURES = [
    "统一 AI 网关（Kimi 大模型，SSE 流式）",
    "KMS 知识库工具（searchKms / getKmsPage）",
    "全局 AI 抽屉「DSTE 智脑」（跨页面上下文感知）",
    "统一 AI 交互 UI（markdown 渲染 · shimmer 思考态）",
]


def set_text(tf, runs_spec, align=PP_ALIGN.LEFT, space_after=Pt(4), line_spacing=1.0):
    """runs_spec: list of paragraphs; each paragraph is list of (text, size, bold, color)."""
    tf.word_wrap = True
    for i, para in enumerate(runs_spec):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.line_spacing = line_spacing
        for text, size, bold, color in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def add_box(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(shape, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    box.shadow.inherit = False
    return box


def build(out_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式

    # 背景
    bg = add_box(slide, 0, 0, prs.slide_width, prs.slide_height, BG, MSO_SHAPE.RECTANGLE)

    # 标题
    title = slide.shapes.add_textbox(Inches(0.45), Inches(0.28), Inches(12.4), Inches(0.85))
    set_text(
        title.text_frame,
        [
            [("经营分析会管理全流程 ", 26, True, INK), ("AI 功能总览", 26, True, RGBColor(0x1F, 0x6F, 0xEB))],
            [("会前准备 · 会中陪伴 · 会后闭环，统一 AI 底座支撑", 12, False, MUTED)],
        ],
        space_after=Pt(2),
    )

    # 三列：会前 / 会中 / 会后
    margin_x = 0.45
    gap = 0.25
    col_w = (13.333 - 2 * margin_x - 2 * gap) / 3
    col_y = 1.35
    col_h = 4.78

    for idx, phase in enumerate(PHASES):
        x = margin_x + idx * (col_w + gap)
        x_emu, y_emu = Inches(x), Inches(col_y)
        w_emu, h_emu = Inches(col_w), Inches(col_h)

        # 卡片底
        add_box(slide, x_emu, y_emu, w_emu, h_emu, CARD_BG)
        # 顶部色条
        add_box(slide, x_emu, y_emu, w_emu, Inches(0.09), phase["color"], MSO_SHAPE.RECTANGLE)

        # 阶段标题
        head = slide.shapes.add_textbox(x_emu + Inches(0.18), y_emu + Inches(0.18), w_emu - Inches(0.36), Inches(0.6))
        set_text(
            head.text_frame,
            [[(phase["name"], 17, True, phase["color"]), ("  " + phase["sub"], 11.5, False, MUTED)]],
        )

        # 功能条目
        body = slide.shapes.add_textbox(x_emu + Inches(0.18), y_emu + Inches(0.85), w_emu - Inches(0.36), h_emu - Inches(1.0))
        paras = []
        for name, desc in phase["items"]:
            paras.append([("▪ ", 11, True, phase["color"]), (name, 11.5, True, INK)])
            paras.append([(desc, 10, False, MUTED)])
        set_text(body.text_frame, paras, space_after=Pt(6), line_spacing=1.08)

        # 阶段间箭头
        if idx < 2:
            arrow = slide.shapes.add_textbox(Inches(x + col_w - 0.02), Inches(col_y + col_h / 2 - 0.2), Inches(gap + 0.06), Inches(0.4))
            set_text(arrow.text_frame, [[("→", 18, True, MUTED)]], align=PP_ALIGN.CENTER)

    # 底部 AI 底座
    base_y = Inches(6.35)
    base = add_box(slide, Inches(margin_x), base_y, Inches(13.333 - 2 * margin_x), Inches(0.88), BASE_BG)
    tf = base.text_frame
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(
        tf,
        [
            [("AI 底座  ", 12.5, True, RGBColor(0x7D, 0xD3, 0xFC)),
             ("    ·    ".join(BASE_FEATURES), 10.5, False, RGBColor(0xE5, 0xE7, 0xEB))],
        ],
        line_spacing=1.15,
    )

    prs.save(out_path)
    print(f"已生成: {out_path}")


if __name__ == "__main__":
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_OUT
    build(out)
