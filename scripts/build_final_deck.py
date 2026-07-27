#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""重建 DSTE汇报合集.pptx（去重后最终版）。

以《经分会事不过三机制-帆软母版版.pptx》为基底（4 页，帆软风格母版），
追加一页帆软风格的《经营分析会管理全流程 AI 功能总览》，输出 5 页终版。

用法:
    .pptx-venv/bin/python scripts/build_final_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
BASE = REPO_ROOT / "经分会事不过三机制-帆软母版版.pptx"
OUT = REPO_ROOT / "DSTE汇报合集.pptx"

FONT = "PingFang SC"

# 帆软母版版配色（从模板中提取）
BLUE = RGBColor(0x03, 0x5D, 0xCF)        # 帆软主蓝
BLUE_MID = RGBColor(0x5B, 0x8C, 0xC8)    # 辅助蓝
CARD_LIGHT = RGBColor(0xE8, 0xF4, 0xFD)  # 浅蓝卡片
CARD_MID = RGBColor(0xBF, 0xDC, 0xFC)    # 中蓝卡片
INK = RGBColor(0x26, 0x26, 0x26)
MUTED = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

PHASES = [
    {
        "name": "会前", "sub": "议程与材料准备", "fill": CARD_LIGHT, "head": BLUE,
        "items": [
            ("AI 议程推荐", "融合历史议程、顺延议题、未闭环行动项/决议与 OMP 重点工作，给出置信度与推荐理由，人工采纳后一键写入议程"),
            ("材料智能审核（KMS）", "4 套场景化评分矩阵，输出总分、维度打分理由、问题清单、改进建议、亮点与审核结论，支持批量审核与历史对比"),
            ("一键 / 批量送审", "议程行直接送审，评分自动回传并显示彩色得分徽标，会议详情页汇总平均分"),
            ("会前准备度检查", "自动计算准备度清单，材料审核评分 ≥60 才算通过"),
        ],
    },
    {
        "name": "会中", "sub": "智能陪伴与执行", "fill": CARD_MID, "head": BLUE,
        "items": [
            ("会议 AI 助手", "注入当前会议上下文（议程/纪要/行动项/决议）的流式问答，附快捷问题芯片，断网降级为本地规则回复"),
            ("工具调用", "实时查询议程、行动项、决议执行情况，边开会边取证"),
            ("AI 草拟 · 人工确认", "AI 生成行动项/新会议草案卡片，点击确认后才写入，可控可追溯"),
        ],
    },
    {
        "name": "会后", "sub": "评估与闭环追踪", "fill": BLUE, "head": WHITE,
        "items": [
            ("AI 自动评分", "三段式模型（会前 35 + 会中 30 + 会后 35），直接消费材料审核分，议程顺延自动扣分，生成反馈标签"),
            ("闭环追踪问答", "一句话查询未闭环行动项、决议执行状态，自动生成经营分析周报草稿"),
        ],
    },
]

BASE_FEATURES = "统一 AI 网关（Kimi 大模型，SSE 流式）  ·  KMS 知识库工具（searchKms / getKmsPage）  ·  全局 AI 抽屉「DSTE 智脑」  ·  统一 AI 交互 UI"


def set_text(tf, paras, align=PP_ALIGN.LEFT, space_after=Pt(4), line_spacing=1.0):
    tf.word_wrap = True
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.line_spacing = line_spacing
        for text, size, bold, color in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def add_box(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(shape, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    box.shadow.inherit = False
    return box


def build_ai_overview_slide(prs):
    blank = next(l for l in prs.slide_masters[0].slide_layouts if l.name == "空白")
    slide = prs.slides.add_slide(blank)

    # 标题区（与母版版内容页一致：0.72 / 0.55 / 1.2）
    title = slide.shapes.add_textbox(Inches(0.72), Inches(0.55), Inches(11.9), Inches(0.7))
    set_text(title.text_frame, [[("经营分析会管理全流程 ", 28, True, INK), ("AI 功能总览", 28, True, BLUE)]])
    sub = slide.shapes.add_textbox(Inches(0.72), Inches(1.2), Inches(11.9), Inches(0.4))
    set_text(sub.text_frame, [[("会前准备 · 会中陪伴 · 会后闭环，统一 AI 底座支撑", 13, False, MUTED)]])

    # 三列卡片：浅蓝 → 中蓝 → 帆软蓝（呼应机制页的三段递进）
    margin_x, gap = 0.72, 0.35
    col_w = (13.333 - 2 * margin_x - 2 * gap) / 3  # ≈ 3.73
    col_y, col_h = 1.85, 4.25

    for idx, phase in enumerate(PHASES):
        x = margin_x + idx * (col_w + gap)
        dark = phase["fill"] == BLUE
        body_color = WHITE if dark else MUTED
        name_color = WHITE if dark else INK

        add_box(slide, Inches(x), Inches(col_y), Inches(col_w), Inches(col_h), phase["fill"])

        head = slide.shapes.add_textbox(Inches(x + 0.22), Inches(col_y + 0.2), Inches(col_w - 0.44), Inches(0.55))
        set_text(head.text_frame, [[(phase["name"], 17, True, phase["head"]),
                                    ("  " + phase["sub"], 11, False, body_color)]])

        body = slide.shapes.add_textbox(Inches(x + 0.22), Inches(col_y + 0.85), Inches(col_w - 0.44), Inches(col_h - 1.05))
        paras = []
        for name, desc in phase["items"]:
            paras.append([("▪ ", 11, True, phase["head"]), (name, 11.5, True, name_color)])
            paras.append([(desc, 9.5, False, body_color)])
        set_text(body.text_frame, paras, space_after=Pt(6), line_spacing=1.08)

        if idx < 2:  # 列间箭头（沿用机制页的 5B8CC8 箭头）
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + col_w + 0.03), Inches(col_y + col_h / 2 - 0.15), Inches(0.29), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = BLUE_MID
            arrow.line.fill.background()
            arrow.shadow.inherit = False

    # 底部 AI 底座条（帆软蓝实底白字）
    bar = add_box(slide, Inches(margin_x), Inches(6.35), Inches(13.333 - 2 * margin_x), Inches(0.75), BLUE)
    tf = bar.text_frame
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(tf, [[("AI 底座   ", 12, True, WHITE), (BASE_FEATURES, 10, False, WHITE)]], line_spacing=1.1)


def main():
    prs = Presentation(str(BASE))
    build_ai_overview_slide(prs)
    prs.save(str(OUT))
    print(f"已生成: {OUT}（共 {len(prs.slides)} 页）")


if __name__ == "__main__":
    main()
