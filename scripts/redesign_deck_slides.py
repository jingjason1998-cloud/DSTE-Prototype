#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""重做 DSTE汇报合集.pptx 的第 5、6 页（pptx-design skill 版，含图标）。

- 第 5 页「6 个问题摸清战略管理骨架」：脊柱图 —— 细线贯穿页面，6 个图标圆点
  节点上下交替排布，幽灵数字做层次；底部 hairline + 一行安静的金句。
- 第 6 页「会议价值对比」：非对称分屏 —— 左侧巨大 90% 数字锚点 + 跨页呼应，
  右侧帆软蓝出血色块 + 出血大纸飞机图标 + chevron 前进符号。

图标由 scripts/build_icon_pngs.mjs 生成（scripts/assets/icons/）。

用法:
    .pptx-venv/bin/python scripts/redesign_deck_slides.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
DECK = REPO_ROOT / "DSTE汇报合集.pptx"
ICON_DIR = REPO_ROOT / "scripts/assets/icons"

FONT = "PingFang SC"

BLUE = RGBColor(0x03, 0x5D, 0xCF)
BLUE_MID = RGBColor(0x5B, 0x8C, 0xC8)
CARD_LIGHT = RGBColor(0xE8, 0xF4, 0xFD)
GHOST = RGBColor(0xD8, 0xE9, 0xFB)
INK = RGBColor(0x26, 0x26, 0x26)
MUTED = RGBColor(0x59, 0x59, 0x59)
HAIRLINE = RGBColor(0xD9, 0xDE, 0xE4)
SPINE = RGBColor(0xC9, 0xDD, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

MARGIN = 0.72
PAGE_W = 13.333
CONTENT_W = PAGE_W - 2 * MARGIN
BRAND = "DSTE 战略管理平台 · 经分会智能闭环"

QUESTIONS = [
    ("01", "大家愿不愿意一起干？", "战略共识", "给所有人装上共同的导航", "users-three"),
    ("02", "目标能不能实现？", "战略解码", "从“要去哪”到“怎么走”", "map-trifold"),
    ("03", "组织能不能顺畅运转？", "流程型组织", "组织跟着战略走", "tree-structure"),
    ("04", "大家有没有动力好好干？", "利益机制", "分钱方向就是用力方向", "currency-dollar"),
    ("05", "日常管理有没有章法？", "会议与督办", "聚焦关键任务", "calendar-check"),
    ("06", "能不能越做越好？", "复盘与迭代", "最大的浪费是经验的浪费", "arrows-clockwise"),
]


def set_text(tf, paras, align=PP_ALIGN.LEFT, space_after=Pt(4), line_spacing=1.0, spacing=None):
    tf.word_wrap = True
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.line_spacing = line_spacing
        for text, size, bold, color in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            if spacing:
                r.font._rPr.set("spc", str(spacing))


def add_shape(slide, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=None, **kw):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor:
        tf.vertical_anchor = anchor
    set_text(tf, paras, align=align, **kw)
    return box


def add_icon(slide, name, color, x, y, size):
    path = ICON_DIR / f"{name}-{color}.png"
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(size), Inches(size))


def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    sldId = list(xml_slides)[index]
    rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rId)
    xml_slides.remove(sldId)


def blank_layout(prs):
    layouts = prs.slide_masters[0].slide_layouts
    return next((l for l in layouts if l.name in ("空白", "Blank")), layouts[6])


# ---------------------------------------------------------------- 第 5 页：脊柱图
def build_spine_slide(prs):
    slide = prs.slides.add_slide(blank_layout(prs))

    add_text(slide, MARGIN, 0.55, 8, 0.3,
             [[("华为 DSTE 系列 · 全景解读", 11, False, MUTED)]], spacing=200)
    add_text(slide, MARGIN, 0.92, 11.9, 0.65,
             [[("6 个问题", 28, True, BLUE), ("，摸清战略管理的骨架", 28, True, INK)]])
    add_text(slide, MARGIN, 1.58, 11.9, 0.35,
             [[("从一张 PPT 到一个持续运转的战略引擎 —— 顺着这根脊梁骨，逐节检查", 13, False, MUTED)]])

    # 脊柱线 + 图标圆点节点
    spine_y = 4.32
    x0, x1 = 0.9, 12.43
    add_shape(slide, x0, spine_y, x1 - x0, 0.018, SPINE)
    col_w = (x1 - x0) / 6

    for i, (num, question, mech, desc, icon) in enumerate(QUESTIONS):
        cx = x0 + (i + 0.5) * col_w
        add_shape(slide, cx - 0.2, spine_y - 0.19, 0.4, 0.4, BLUE, MSO_SHAPE.OVAL)
        add_icon(slide, icon, "FFFFFF", cx - 0.11, spine_y - 0.10, 0.22)

        tx = cx - 0.93
        tw = 1.86
        if i % 2 == 0:  # 上方
            add_text(slide, tx, 2.02, tw, 0.62, [[(num, 40, True, GHOST)]], align=PP_ALIGN.CENTER)
            add_text(slide, tx, 2.72, tw, 0.55, [[(question, 12, True, INK)]],
                     align=PP_ALIGN.CENTER, line_spacing=1.1)
            add_text(slide, tx, 3.42, tw, 0.28, [[(mech, 10, True, BLUE)]], align=PP_ALIGN.CENTER)
            add_text(slide, tx, 3.68, tw, 0.42, [[(desc, 8.5, False, MUTED)]],
                     align=PP_ALIGN.CENTER, line_spacing=1.15)
        else:  # 下方
            add_text(slide, tx, 4.72, tw, 0.62, [[(num, 40, True, GHOST)]], align=PP_ALIGN.CENTER)
            add_text(slide, tx, 5.42, tw, 0.55, [[(question, 12, True, INK)]],
                     align=PP_ALIGN.CENTER, line_spacing=1.1)
            add_text(slide, tx, 6.12, tw, 0.28, [[(mech, 10, True, BLUE)]], align=PP_ALIGN.CENTER)
            add_text(slide, tx, 6.38, tw, 0.42, [[(desc, 8.5, False, MUTED)]],
                     align=PP_ALIGN.CENTER, line_spacing=1.15)

    # 底部：hairline + 左品牌 + 右金句页码
    add_shape(slide, MARGIN, 6.98, CONTENT_W, 0.012, HAIRLINE)
    add_text(slide, MARGIN, 7.08, 6, 0.3, [[(BRAND, 9, False, MUTED)]])
    add_text(slide, PAGE_W - MARGIN - 7, 7.08, 7, 0.3,
             [[("认准一套不复杂的方法，几十年死磕到底 —— 大道至简　05 / 06", 9.5, False, MUTED)]],
             align=PP_ALIGN.RIGHT)


# -------------------------------------------------------------- 第 6 页：非对称分屏
def build_split_slide(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    panel_x = 8.2

    # 右侧帆软蓝出血色块（贴边，顶到底）
    add_shape(slide, panel_x, 0, PAGE_W - panel_x, 7.5, BLUE)
    # 出血大纸飞机（浅蓝幽灵图标，破格压底边）
    add_icon(slide, "paper-plane-right", "BFDCFC", 10.85, 5.0, 2.6)

    # 左侧：eyebrow → 巨大数字锚点 → 结论 → 跨页呼应
    add_icon(slide, "users-three", "B8BCC2", MARGIN, 0.52, 0.24)
    add_text(slide, MARGIN + 0.34, 0.5, 7, 0.3,
             [[("华为“开会”的三个规则", 11, False, MUTED)]], spacing=200)
    add_text(slide, 0.68, 1.3, 5, 1.5, [[("90%", 84, True, BLUE)]])
    add_text(slide, MARGIN, 3.0, 7.2, 0.55,
             [[("的企业，会议只是", 26, True, INK), ("社交聚集", 26, True, INK)]])
    add_text(slide, MARGIN, 3.75, 7.2, 0.35,
             [[("消耗时间，原地踏步 —— 为什么绝大多数企业做不到", 12, False, MUTED)]])

    add_shape(slide, MARGIN, 5.9, 3.2, 0.012, HAIRLINE)
    add_text(slide, MARGIN, 6.08, 7.2, 0.6,
             [[("DSTE 的解法：", 10.5, False, MUTED),
               ("经分会「事不过三」议程顺延机制", 10.5, True, BLUE),
               ("，见本合集第 1 页", 10.5, False, MUTED)]], line_spacing=1.3)
    add_text(slide, MARGIN, 7.08, 7, 0.3, [[(BRAND + "　06 / 06", 9, False, MUTED)]])

    # 右侧蓝块内：华为 = 推进引擎
    px_ = panel_x + 0.55
    add_text(slide, px_, 1.5, 4.2, 0.3,
             [[("华为 HUAWEI", 12, True, CARD_LIGHT)]], spacing=300)
    add_text(slide, px_, 2.1, 4.3, 0.55, [[("会议 = 推进引擎", 27, True, WHITE)]])
    add_shape(slide, px_, 3.08, 1.2, 0.014, BLUE_MID)
    add_text(slide, px_, 3.32, 4.2, 0.7,
             [[("几乎每一场会议，", 12.5, False, CARD_LIGHT)],
              [("都在推着事情往前走", 12.5, False, CARD_LIGHT)]], line_spacing=1.35)

    # 三个 chevron：前进的动势
    for j in range(3):
        add_shape(slide, px_ + j * 0.62, 4.5, 0.52, 0.46,
                  WHITE if j == 2 else (CARD_LIGHT if j == 1 else BLUE_MID),
                  MSO_SHAPE.CHEVRON)


def main():
    prs = Presentation(str(DECK))
    # 删除上一版的最后两页
    delete_slide(prs, len(prs.slides._sldIdLst) - 1)
    delete_slide(prs, len(prs.slides._sldIdLst) - 1)
    build_spine_slide(prs)
    build_split_slide(prs)
    prs.save(str(DECK))
    print(f"已重做最后两页: {DECK}（共 {len(prs.slides._sldIdLst)} 页）")


if __name__ == "__main__":
    main()
